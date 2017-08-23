import json
import sys
import os.path

from pptx import Presentation
from pptx.util import Inches
from six.moves import urllib

# {
#     "image_url": "https://pbs.twimg.com/profile_images/894625939652579328/Cmbsq0OP_400x400.jpg",
#     "text": "OpenFaaS",
#     "title": "Serverless on your servers.. yas"
# }

def maybe_download_url(url, dest_directory):
    if not os.path.exists(dest_directory):
        os.makedirs(dest_directory)

    filename = url.split('/')[-1]
    filepath = os.path.join(dest_directory, filename)
    if not os.path.exists(filepath):
        filepathResult, _ = urllib.request.urlretrieve(url, filepath)

        # statinfo = os.stat(filepath)

    return filepath

def make_prezo(req):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]

    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title

    subtitle = slide.placeholders[1]
    title.text = req["title"]
    subtitle.text = req["text"]

    top = Inches(1)
    left = Inches(2)
    height = Inches(1.5)

    img = maybe_download_url(req["image_url"], "/tmp")
    pic = slide.shapes.add_picture(img, left, top, height=height)
    return prs

def output_prezo(prs):
    prs.save("out.pptx")
    myfile = open('out.pptx', 'r')
    data = myfile.read()
    myfile.close()
    sys.stdout.write(data)


def handle(val):
    req = json.loads(val)

    prs = make_prezo(req)

    output_prezo(prs)
